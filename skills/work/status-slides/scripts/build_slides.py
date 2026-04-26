#!/usr/bin/env python3
"""
build_slides.py — Convert a team-status markdown page into a best-practice
executive PowerPoint deck.

Canonical 7-slide structure (audience-conditional inclusion):
  1. Title
  2. Executive Summary  (## Synopsis)
  3. RAG Dashboard      (## Progress — markdown table)
  4. Top Risks          (## Risks)
  5. Open Issues        (## Issues — dropped for audience=customer)
  6. Asks / Decisions   (## Asks)
  7. What's Next        (## What's Next)

Usage:
  python build_slides.py \
    --source wiki/projects/order-platform/status/2026-04-25-team-status.md \
    --output outputs/2026-04-25-order-platform-status.pptx \
    --theme default \
    --audience leadership

Requires: pip install python-pptx
"""
from __future__ import annotations

import argparse
import re
import sys
from datetime import date
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print(
        "ERROR: python-pptx not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)


# ---------- themes ----------

THEMES = {
    "default": {
        "title_pt": 36,
        "heading_pt": 28,
        "body_pt": 18,
        "table_pt": 14,
        "primary": RGBColor(0x1F, 0x3A, 0x5F),
        "accent": RGBColor(0xC4, 0x4D, 0x58),
        "muted": RGBColor(0x66, 0x66, 0x66),
        "rag_green": RGBColor(0x2C, 0x8A, 0x4D),
        "rag_amber": RGBColor(0xE0, 0x9F, 0x3E),
        "rag_red": RGBColor(0xC4, 0x33, 0x33),
    },
    "minimal": {
        "title_pt": 32,
        "heading_pt": 24,
        "body_pt": 16,
        "table_pt": 12,
        "primary": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x00, 0x66, 0xCC),
        "muted": RGBColor(0x88, 0x88, 0x88),
        "rag_green": RGBColor(0x33, 0x99, 0x33),
        "rag_amber": RGBColor(0xCC, 0x99, 0x33),
        "rag_red": RGBColor(0xCC, 0x33, 0x33),
    },
    "technical": {
        "title_pt": 32,
        "heading_pt": 22,
        "body_pt": 14,
        "table_pt": 11,
        "primary": RGBColor(0x00, 0x33, 0x66),
        "accent": RGBColor(0x00, 0x99, 0x99),
        "muted": RGBColor(0x55, 0x55, 0x55),
        "rag_green": RGBColor(0x2C, 0x8A, 0x4D),
        "rag_amber": RGBColor(0xE0, 0x9F, 0x3E),
        "rag_red": RGBColor(0xC4, 0x33, 0x33),
    },
}


# ---------- markdown parsing ----------

def parse_frontmatter(text: str):
    m = re.match(r"^---\n(.*?)\n---\n(.*)$", text, re.DOTALL)
    if not m:
        return {}, text
    fm: dict[str, str] = {}
    for line in m.group(1).splitlines():
        if ":" in line and not line.startswith(" "):
            k, _, v = line.partition(":")
            fm[k.strip()] = v.strip().strip('"')
    return fm, m.group(2)


def split_sections(body: str):
    """Split markdown body into {h2-heading: content}."""
    sections: dict[str, str] = {}
    current = None
    buf: list[str] = []
    for line in body.splitlines():
        m = re.match(r"^##\s+(.+?)\s*$", line)
        if m:
            if current is not None:
                sections[current] = "\n".join(buf).strip()
            current = m.group(1).strip()
            buf = []
        else:
            buf.append(line)
    if current is not None:
        sections[current] = "\n".join(buf).strip()
    return sections


def section_lookup(sections: dict[str, str], *names: str) -> str:
    """Find first section whose key matches any of the names (case-insensitive, partial)."""
    for name in names:
        n = name.lower()
        for k in sections:
            kl = k.lower()
            if kl == n or n in kl:
                return sections[k]
    return ""


def parse_table(text: str):
    """Parse a markdown table; return list[dict[header→cell]] or []."""
    lines = [l for l in text.splitlines() if l.strip().startswith("|")]
    if len(lines) < 2:
        return []
    headers = [h.strip() for h in lines[0].strip("|").split("|")]
    rows = []
    for line in lines[2:]:  # skip the separator row
        cells = [c.strip() for c in line.strip("|").split("|")]
        if len(cells) == len(headers):
            rows.append(dict(zip(headers, cells)))
    return rows


def parse_bullets(text: str):
    bullets: list[str] = []
    for line in text.splitlines():
        m = re.match(r"^\s*[-*]\s+(.+)$", line)
        if m:
            bullets.append(m.group(1).strip())
    return bullets


# ---------- slide builders ----------

def add_title_slide(prs, theme, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = tx.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(theme["title_pt"])
    p.font.bold = True
    p.font.color.rgb = theme["primary"]
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.6))
        sub.text_frame.text = subtitle
        sp = sub.text_frame.paragraphs[0]
        sp.font.size = Pt(theme["body_pt"])
        sp.font.color.rgb = theme["muted"]


def add_bullet_slide(prs, theme, title, bullets, max_bullets=8):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    tp = title_box.text_frame.paragraphs[0]
    tp.font.size = Pt(theme["heading_pt"])
    tp.font.bold = True
    tp.font.color.rgb = theme["primary"]

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    items = bullets[:max_bullets] if bullets else ["[FILL IN]"]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(theme["body_pt"])
        p.font.color.rgb = theme["accent"] if item == "[FILL IN]" else theme["primary"]
    if bullets and len(bullets) > max_bullets:
        p = tf.add_paragraph()
        p.text = f"  (+{len(bullets) - max_bullets} more — see appendix)"
        p.font.size = Pt(max(theme["body_pt"] - 2, 10))
        p.font.italic = True
        p.font.color.rgb = theme["muted"]


def add_table_slide(prs, theme, title, rows, headers):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    tp = title_box.text_frame.paragraphs[0]
    tp.font.size = Pt(theme["heading_pt"])
    tp.font.bold = True
    tp.font.color.rgb = theme["primary"]

    if not rows:
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        body.text_frame.text = "[FILL IN — Progress table missing in source]"
        bp = body.text_frame.paragraphs[0]
        bp.font.size = Pt(theme["body_pt"])
        bp.font.color.rgb = theme["accent"]
        return

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 header
    height = Inches(min(0.5 * n_rows, 5.5))
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(9), height
    ).table

    # header
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(theme["table_pt"])
            p.font.bold = True
            p.font.color.rgb = theme["primary"]

    # data rows
    for i, row in enumerate(rows, start=1):
        for j, h in enumerate(headers):
            c = tbl.cell(i, j)
            val = row.get(h, "")
            c.text = val
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(theme["table_pt"])
                lower = val.lower()
                if "🟢" in val or " green" in lower or lower.strip() == "green":
                    p.font.color.rgb = theme["rag_green"]
                elif "🟡" in val or " amber" in lower or lower.strip() == "amber":
                    p.font.color.rgb = theme["rag_amber"]
                elif "🔴" in val or " red" in lower or lower.strip() == "red":
                    p.font.color.rgb = theme["rag_red"]
                else:
                    p.font.color.rgb = theme["primary"]


# ---------- main ----------

def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(
        description="Build a status PowerPoint deck from a markdown status page."
    )
    ap.add_argument("--source", required=True, help="Path to status markdown")
    ap.add_argument("--output", required=True, help="Output .pptx path")
    ap.add_argument("--theme", choices=THEMES.keys(), default="default")
    ap.add_argument(
        "--audience",
        choices=["internal", "leadership", "steering-committee", "customer"],
        default="leadership",
    )
    ap.add_argument(
        "--template",
        default=None,
        help="Optional brand-styled .pptx template to start from",
    )
    args = ap.parse_args(argv)

    src = Path(args.source)
    if not src.exists():
        print(f"ERROR: source not found: {src}", file=sys.stderr)
        return 1

    text = src.read_text()
    fm, body = parse_frontmatter(text)
    sections = split_sections(body)
    theme = THEMES[args.theme]

    prs = Presentation(args.template) if args.template else Presentation()

    # 1. Title
    project = fm.get("project") or src.parent.parent.name
    period = fm.get("period", str(date.today()))
    audience = fm.get("audience", args.audience)
    add_title_slide(
        prs, theme, f"{project} — Team Status", f"{period} · {audience}"
    )

    # 2. Executive summary
    syn = section_lookup(sections, "Synopsis", "Executive Summary", "Summary")
    syn_bullets = parse_bullets(syn) or ([syn.strip()] if syn.strip() else [])
    add_bullet_slide(prs, theme, "Executive Summary", syn_bullets, max_bullets=4)

    # 3. RAG dashboard
    progress_text = section_lookup(sections, "Progress")
    rows = parse_table(progress_text)
    if rows:
        add_table_slide(prs, theme, "RAG Dashboard", rows, list(rows[0].keys()))
    else:
        add_bullet_slide(prs, theme, "Progress", parse_bullets(progress_text), max_bullets=8)

    # 4. Top risks (cap by audience)
    cap = 5 if args.audience in ("leadership", "steering-committee") else 8
    add_bullet_slide(
        prs, theme, "Top Risks", parse_bullets(section_lookup(sections, "Risks")), max_bullets=cap
    )

    # 5. Open issues — dropped for customer audience
    if args.audience != "customer":
        add_bullet_slide(
            prs, theme, "Open Issues",
            parse_bullets(section_lookup(sections, "Issues")), max_bullets=cap,
        )

    # 6. Asks / Decisions needed
    asks = parse_bullets(section_lookup(sections, "Asks", "Decisions Needed"))
    add_bullet_slide(prs, theme, "Asks / Decisions Needed", asks, max_bullets=cap)

    # 7. What's next
    nxt = parse_bullets(section_lookup(sections, "What's Next", "Next Steps", "Next"))
    add_bullet_slide(prs, theme, "What's Next", nxt, max_bullets=8)

    out = Path(args.output)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(
        f"Wrote {out} ({len(prs.slides)} slides, theme={args.theme}, audience={args.audience})"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
